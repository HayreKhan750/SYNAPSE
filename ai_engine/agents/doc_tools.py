"""
ai_engine.agents.doc_tools
~~~~~~~~~~~~~~~~~~~~~~~~~~
Document generation tools for SYNAPSE ReAct agents.

Phase 5.2 — Document Generation (Week 14)

Tools implemented:
  1. generate_pdf       — ReportLab styled PDF with cover, TOC, sections, footer
  2. generate_ppt       — python-pptx PowerPoint with title + content slides
  3. generate_word_doc  — python-docx Word document with styles, headings, TOC
  4. generate_markdown  — Plain markdown document

Each tool:
  - Saves the file to MEDIA_ROOT/documents/<user_id>/<uuid>.<ext>
  - Returns a plain string with file path + metadata (agent-readable)
  - Handles errors gracefully (returns error string, never raises)
  - Is registered via make_*_tool() factory functions
"""
from __future__ import annotations

import logging
import os
import uuid
from pathlib import Path
from typing import Any, Dict, List, Optional

from langchain_core.tools import StructuredTool
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Storage root helper
# ---------------------------------------------------------------------------

def _doc_dir(user_id: str = "anonymous") -> Path:
    """Return (and create) the per-user document storage directory."""
    # Prefer DJANGO_MEDIA_ROOT (set in .env), fall back to MEDIA_ROOT, then 'media'
    media_root = Path(
        os.environ.get("DJANGO_MEDIA_ROOT") or
        os.environ.get("MEDIA_ROOT") or
        "media"
    )
    doc_dir = media_root / "documents" / str(user_id)
    doc_dir.mkdir(parents=True, exist_ok=True)
    return doc_dir


def _rel_path(abs_path: Path) -> str:
    """Return a media-relative path string for storing in the DB."""
    media_root = Path(
        os.environ.get("DJANGO_MEDIA_ROOT") or
        os.environ.get("MEDIA_ROOT") or
        "media"
    )
    try:
        return str(abs_path.relative_to(media_root))
    except ValueError:
        return str(abs_path)


# ===========================================================================
# 1. generate_pdf
# ===========================================================================

class GeneratePDFInput(BaseModel):
    title: str = Field(..., description="Document title shown on the cover page")
    sections: List[Dict[str, str]] = Field(
        ...,
        description=(
            "List of sections. Each section is a dict with 'heading' (str) and 'content' (str). "
            "Example: [{'heading': 'Introduction', 'content': 'Text here...'}]"
        ),
    )
    subtitle: str = Field(default="", description="Optional subtitle shown on the cover page")
    author: str = Field(default="SYNAPSE AI", description="Author name shown in footer")
    user_id: str = Field(default="anonymous", description="User ID for file storage path")


def _generate_pdf(
    title: str,
    sections: List[Dict[str, str]],
    subtitle: str = "",
    author: str = "SYNAPSE AI",
    user_id: str = "anonymous",
) -> str:
    """Generate a styled PDF report using ReportLab."""
    try:
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import (
            HRFlowable, PageBreak, Paragraph, SimpleDocTemplate, Spacer,
        )

        file_name = f"{uuid.uuid4().hex}.pdf"
        file_path = _doc_dir(user_id) / file_name

        # ── Styles ────────────────────────────────────────────────────
        styles = getSampleStyleSheet()

        style_cover_title = ParagraphStyle(
            "CoverTitle",
            parent=styles["Title"],
            fontSize=32,
            textColor=colors.HexColor("#4F46E5"),
            spaceAfter=12,
            alignment=TA_CENTER,
            fontName="Helvetica-Bold",
        )
        style_cover_subtitle = ParagraphStyle(
            "CoverSubtitle",
            parent=styles["Normal"],
            fontSize=16,
            textColor=colors.HexColor("#6B7280"),
            spaceAfter=8,
            alignment=TA_CENTER,
        )
        style_cover_author = ParagraphStyle(
            "CoverAuthor",
            parent=styles["Normal"],
            fontSize=12,
            textColor=colors.HexColor("#9CA3AF"),
            alignment=TA_CENTER,
        )
        style_h1 = ParagraphStyle(
            "SynapseH1",
            parent=styles["Heading1"],
            fontSize=18,
            textColor=colors.HexColor("#1E1B4B"),
            spaceBefore=18,
            spaceAfter=8,
            fontName="Helvetica-Bold",
        )
        style_body = ParagraphStyle(
            "SynapseBody",
            parent=styles["Normal"],
            fontSize=11,
            leading=16,
            textColor=colors.HexColor("#374151"),
            spaceAfter=10,
            alignment=TA_JUSTIFY,
        )

        # ── Page template with footer ──────────────────────────────────
        def _on_page(canvas, doc):
            canvas.saveState()
            canvas.setFont("Helvetica", 8)
            canvas.setFillColor(colors.HexColor("#9CA3AF"))
            canvas.drawString(2 * cm, 1 * cm, f"{title} — {author}")
            canvas.drawRightString(A4[0] - 2 * cm, 1 * cm, f"Page {doc.page}")
            canvas.restoreState()

        doc = SimpleDocTemplate(
            str(file_path),
            pagesize=A4,
            rightMargin=2 * cm,
            leftMargin=2 * cm,
            topMargin=2.5 * cm,
            bottomMargin=2 * cm,
        )

        story = []

        # ── Cover page ────────────────────────────────────────────────
        story.append(Spacer(1, 4 * cm))
        story.append(Paragraph(title, style_cover_title))
        if subtitle:
            story.append(Paragraph(subtitle, style_cover_subtitle))
        story.append(Spacer(1, 0.5 * cm))
        story.append(HRFlowable(width="80%", thickness=2, color=colors.HexColor("#4F46E5")))
        story.append(Spacer(1, 0.5 * cm))
        story.append(Paragraph(f"Generated by {author}", style_cover_author))
        story.append(PageBreak())

        # ── Sections ─────────────────────────────────────────────────
        for i, section in enumerate(sections):
            heading = section.get("heading", f"Section {i + 1}")
            content = section.get("content", "")

            story.append(Paragraph(heading, style_h1))
            story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor("#E5E7EB")))
            story.append(Spacer(1, 0.3 * cm))

            # Split content into paragraphs on double newline
            for para in content.split("\n\n"):
                para = para.strip()
                if para:
                    # Escape XML special chars for ReportLab
                    para_safe = para.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    story.append(Paragraph(para_safe, style_body))

            story.append(Spacer(1, 0.5 * cm))

        doc.build(story, onFirstPage=_on_page, onLaterPages=_on_page)

        file_size = file_path.stat().st_size
        rel = _rel_path(file_path)

        return (
            f"PDF generated successfully.\n"
            f"Title: {title}\n"
            f"Sections: {len(sections)}\n"
            f"File: {rel}\n"
            f"Size: {file_size:,} bytes\n"
            f"Path: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("generate_pdf failed: %s", exc)
        return f"PDF generation failed: {exc}"


def make_generate_pdf_tool() -> StructuredTool:
    return StructuredTool.from_function(
        func=_generate_pdf,
        name="generate_pdf",
        description=(
            "Generate a styled PDF report with a cover page, section headings, and body text. "
            "Provide a title, optional subtitle, author name, and a list of sections "
            "(each with 'heading' and 'content'). Returns the file path of the saved PDF."
        ),
        args_schema=GeneratePDFInput,
        return_direct=False,
    )


# ===========================================================================
# 2. generate_ppt
# ===========================================================================

class SlideSpec(BaseModel):
    title: str = Field(..., description="Slide title")
    bullets: List[str] = Field(default_factory=list, description="Bullet point lines for the slide")
    notes: str = Field(default="", description="Speaker notes for the slide")


class GeneratePPTInput(BaseModel):
    title: str = Field(..., description="Presentation title (shown on the title slide)")
    subtitle: str = Field(default="Generated by SYNAPSE AI", description="Subtitle on the title slide")
    slides: List[SlideSpec] = Field(..., description="List of content slides")
    author: str = Field(default="SYNAPSE AI", description="Author shown in presentation metadata")
    user_id: str = Field(default="anonymous", description="User ID for file storage path")


def _generate_ppt(
    title: str,
    slides: List[Dict[str, Any]],
    subtitle: str = "Generated by SYNAPSE AI",
    author: str = "SYNAPSE AI",
    user_id: str = "anonymous",
) -> str:
    """Generate a styled PowerPoint presentation using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        file_name = f"{uuid.uuid4().hex}.pptx"
        file_path = _doc_dir(user_id) / file_name

        prs = Presentation()
        prs.core_properties.author = author
        prs.core_properties.title = title

        # SYNAPSE brand colours
        INDIGO = RGBColor(0x4F, 0x46, 0xE5)
        DARK   = RGBColor(0x1E, 0x1B, 0x4B)
        GRAY   = RGBColor(0x6B, 0x72, 0x80)
        WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

        slide_width  = prs.slide_width
        slide_height = prs.slide_height

        def _set_bg(slide, colour: RGBColor = WHITE):
            from pptx.oxml.ns import qn
            from lxml import etree
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = colour

        def _add_text_box(slide, text, left, top, width, height,
                          font_size=18, bold=False, colour=DARK, align=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = colour
            return txBox

        # ── Title slide ───────────────────────────────────────────────
        blank_layout = prs.slide_layouts[6]  # blank layout
        title_slide = prs.slides.add_slide(blank_layout)
        _set_bg(title_slide, INDIGO)

        # Title text
        _add_text_box(
            title_slide, title,
            left=Inches(0.8), top=Inches(2.5),
            width=Inches(8.4), height=Inches(1.2),
            font_size=36, bold=True, colour=WHITE, align=PP_ALIGN.CENTER,
        )
        # Subtitle text
        _add_text_box(
            title_slide, subtitle,
            left=Inches(0.8), top=Inches(3.9),
            width=Inches(8.4), height=Inches(0.7),
            font_size=18, bold=False, colour=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER,
        )
        # Author / footer
        _add_text_box(
            title_slide, author,
            left=Inches(0.8), top=Inches(6.5),
            width=Inches(8.4), height=Inches(0.4),
            font_size=11, colour=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER,
        )

        # ── Content slides ────────────────────────────────────────────
        for slide_data in slides:
            if isinstance(slide_data, dict):
                s_title   = slide_data.get("title", "Slide")
                s_bullets = slide_data.get("bullets", [])
                s_notes   = slide_data.get("notes", "")
            else:
                s_title   = getattr(slide_data, "title", "Slide")
                s_bullets = getattr(slide_data, "bullets", [])
                s_notes   = getattr(slide_data, "notes", "")

            content_slide = prs.slides.add_slide(blank_layout)
            _set_bg(content_slide, WHITE)

            # Title bar (indigo strip at top)
            from pptx.util import Emu
            bar = content_slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                left=0, top=0,
                width=slide_width, height=Inches(1.1),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = INDIGO
            bar.line.fill.background()

            # Slide title
            _add_text_box(
                content_slide, s_title,
                left=Inches(0.4), top=Inches(0.15),
                width=Inches(9.2), height=Inches(0.8),
                font_size=22, bold=True, colour=WHITE,
            )

            # Bullet points
            if s_bullets:
                from pptx.util import Pt as _Pt
                body_box = content_slide.shapes.add_textbox(
                    Inches(0.6), Inches(1.3),
                    Inches(8.8), Inches(5.0),
                )
                tf = body_box.text_frame
                tf.word_wrap = True
                for j, bullet in enumerate(s_bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = f"• {bullet}"
                    p.font.size = _Pt(16)
                    p.font.color.rgb = DARK
                    p.space_after = _Pt(6)

            # Speaker notes
            if s_notes:
                notes_slide = content_slide.notes_slide
                notes_slide.notes_text_frame.text = s_notes

        # ── Slide numbers (footer shape) on each content slide ────────
        for i, slide in enumerate(prs.slides):
            if i == 0:
                continue  # skip title slide
            _add_text_box(
                slide, f"{i}/{len(prs.slides) - 1}",
                left=Inches(8.8), top=Inches(6.8),
                width=Inches(1.0), height=Inches(0.3),
                font_size=9, colour=GRAY, align=PP_ALIGN.RIGHT,
            )

        prs.save(str(file_path))

        file_size = file_path.stat().st_size
        rel = _rel_path(file_path)

        return (
            f"PowerPoint generated successfully.\n"
            f"Title: {title}\n"
            f"Slides: {len(slides) + 1} (including title slide)\n"
            f"File: {rel}\n"
            f"Size: {file_size:,} bytes\n"
            f"Path: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("generate_ppt failed: %s", exc)
        return f"PowerPoint generation failed: {exc}"


def make_generate_ppt_tool() -> StructuredTool:
    return StructuredTool.from_function(
        func=_generate_ppt,
        name="generate_ppt",
        description=(
            "Generate a styled PowerPoint presentation with a branded title slide and content slides. "
            "Provide a title, subtitle, and a list of slides (each with 'title', 'bullets' list, "
            "and optional 'notes'). Returns the file path of the saved .pptx file."
        ),
        args_schema=GeneratePPTInput,
        return_direct=False,
    )


# ===========================================================================
# 3. generate_word_doc
# ===========================================================================

class WordSection(BaseModel):
    heading: str = Field(..., description="Section heading text")
    content: str = Field(..., description="Section body content (paragraphs separated by newlines)")
    level: int = Field(default=1, ge=1, le=3, description="Heading level: 1 = H1, 2 = H2, 3 = H3")


class GenerateWordDocInput(BaseModel):
    title: str = Field(..., description="Document title")
    sections: List[WordSection] = Field(..., description="List of document sections")
    author: str = Field(default="SYNAPSE AI", description="Author metadata")
    add_toc: bool = Field(default=True, description="Whether to add a Table of Contents placeholder")
    user_id: str = Field(default="anonymous", description="User ID for file storage path")


def _generate_word_doc(
    title: str,
    sections: List[Dict[str, Any]],
    author: str = "SYNAPSE AI",
    add_toc: bool = True,
    user_id: str = "anonymous",
) -> str:
    """Generate a styled Word document using python-docx."""
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.shared import Inches, Pt, RGBColor
        from lxml import etree

        file_name = f"{uuid.uuid4().hex}.docx"
        file_path = _doc_dir(user_id) / file_name

        doc = Document()

        # ── Core properties ───────────────────────────────────────────
        doc.core_properties.author = author
        doc.core_properties.title = title

        # ── Page margins ──────────────────────────────────────────────
        for section in doc.sections:
            section.top_margin    = Inches(1.0)
            section.bottom_margin = Inches(1.0)
            section.left_margin   = Inches(1.2)
            section.right_margin  = Inches(1.2)

        # ── Custom styles ─────────────────────────────────────────────
        def _apply_heading_style(paragraph, level: int = 1):
            """Apply SYNAPSE-branded heading colours."""
            from docx.oxml.ns import qn as _qn
            colours = {1: "4F46E5", 2: "1E1B4B", 3: "374151"}
            colour = colours.get(level, "374151")
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor.from_string(colour)
                run.font.bold = True

        # ── Title ─────────────────────────────────────────────────────
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in title_para.runs:
            run.font.color.rgb = RGBColor(0x4F, 0x46, 0xE5)
            run.font.size = Pt(24)

        doc.add_paragraph()  # spacer

        # ── Table of Contents placeholder ────────────────────────────
        if add_toc:
            toc_heading = doc.add_heading("Table of Contents", level=1)
            _apply_heading_style(toc_heading, 1)

            # Insert a proper Word TOC field
            toc_para = doc.add_paragraph()
            fldChar_begin = etree.SubElement(
                etree.SubElement(
                    etree.SubElement(toc_para._p, qn("w:r")),
                    qn("w:fldChar"),
                ),
                qn("w:fldCharType"),
            )
            fldChar_begin.set(qn("w:fldCharType"), "begin")

            run_instr = toc_para._p.add_run() if hasattr(toc_para._p, "add_run") else etree.SubElement(toc_para._p, qn("w:r"))
            instrText = etree.SubElement(run_instr, qn("w:instrText"))
            instrText.text = ' TOC \\o "1-3" \\h \\z \\u '

            fldChar_sep = etree.SubElement(
                etree.SubElement(toc_para._p, qn("w:r")),
                qn("w:fldChar"),
            )
            fldChar_sep.set(qn("w:fldCharType"), "separate")

            toc_text_run = etree.SubElement(toc_para._p, qn("w:r"))
            t = etree.SubElement(toc_text_run, qn("w:t"))
            t.text = "[Right-click and Update Field to generate Table of Contents]"

            fldChar_end = etree.SubElement(
                etree.SubElement(toc_para._p, qn("w:r")),
                qn("w:fldChar"),
            )
            fldChar_end.set(qn("w:fldCharType"), "end")

            doc.add_page_break()

        # ── Sections ─────────────────────────────────────────────────
        for section_data in sections:
            if isinstance(section_data, dict):
                heading_text = section_data.get("heading", "Section")
                content_text = section_data.get("content", "")
                level        = int(section_data.get("level", 1))
            else:
                heading_text = getattr(section_data, "heading", "Section")
                content_text = getattr(section_data, "content", "")
                level        = int(getattr(section_data, "level", 1))

            level = max(1, min(3, level))
            h = doc.add_heading(heading_text, level=level)
            _apply_heading_style(h, level)

            for para_text in content_text.split("\n\n"):
                para_text = para_text.strip()
                if not para_text:
                    continue
                p = doc.add_paragraph(para_text)
                p.paragraph_format.space_after = Pt(8)
                for run in p.runs:
                    run.font.size = Pt(11)

        # ── Footer with page numbers ──────────────────────────────────
        from docx.oxml.ns import qn as _qn
        for sec in doc.sections:
            footer = sec.footer
            footer_para = footer.paragraphs[0]
            footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            footer_para.clear()
            run = footer_para.add_run(f"{title} — {author}  |  Page ")
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)

            # Page number field
            fldChar1 = etree.SubElement(etree.SubElement(footer_para._p, _qn("w:r")), _qn("w:fldChar"))
            fldChar1.set(_qn("w:fldCharType"), "begin")
            instrRun = etree.SubElement(footer_para._p, _qn("w:r"))
            instr = etree.SubElement(instrRun, _qn("w:instrText"))
            instr.text = " PAGE "
            fldChar2 = etree.SubElement(etree.SubElement(footer_para._p, _qn("w:r")), _qn("w:fldChar"))
            fldChar2.set(_qn("w:fldCharType"), "end")

        doc.save(str(file_path))

        file_size = file_path.stat().st_size
        rel = _rel_path(file_path)

        return (
            f"Word document generated successfully.\n"
            f"Title: {title}\n"
            f"Sections: {len(sections)}\n"
            f"TOC: {'yes' if add_toc else 'no'}\n"
            f"File: {rel}\n"
            f"Size: {file_size:,} bytes\n"
            f"Path: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("generate_word_doc failed: %s", exc)
        return f"Word document generation failed: {exc}"


def make_generate_word_doc_tool() -> StructuredTool:
    return StructuredTool.from_function(
        func=_generate_word_doc,
        name="generate_word_doc",
        description=(
            "Generate a styled Microsoft Word document (.docx) with a title, table of contents, "
            "and structured sections. Each section has a 'heading', 'content', and optional 'level' "
            "(1=H1, 2=H2, 3=H3). Returns the file path of the saved .docx file."
        ),
        args_schema=GenerateWordDocInput,
        return_direct=False,
    )


# ===========================================================================
# 4. generate_markdown
# ===========================================================================

class GenerateMarkdownInput(BaseModel):
    title: str = Field(..., description="Document title (becomes H1 heading)")
    sections: List[Dict[str, str]] = Field(
        ...,
        description="List of sections. Each dict has 'heading' and 'content'.",
    )
    author: str = Field(default="SYNAPSE AI", description="Author shown in document header")
    user_id: str = Field(default="anonymous", description="User ID for file storage path")


def _generate_markdown(
    title: str,
    sections: List[Dict[str, str]],
    author: str = "SYNAPSE AI",
    user_id: str = "anonymous",
) -> str:
    """Generate a Markdown document and save to disk."""
    try:
        from datetime import datetime

        file_name = f"{uuid.uuid4().hex}.md"
        file_path = _doc_dir(user_id) / file_name

        lines = [
            f"# {title}",
            f"",
            f"> Generated by {author} on {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}",
            f"",
            f"---",
            f"",
        ]

        # Table of contents
        lines.append("## Table of Contents")
        lines.append("")
        for i, section in enumerate(sections, 1):
            heading = section.get("heading", f"Section {i}")
            anchor = heading.lower().replace(" ", "-").replace("/", "")
            lines.append(f"{i}. [{heading}](#{anchor})")
        lines.append("")
        lines.append("---")
        lines.append("")

        # Sections
        for section in sections:
            heading = section.get("heading", "Section")
            content = section.get("content", "")
            lines.append(f"## {heading}")
            lines.append("")
            lines.append(content.strip())
            lines.append("")
            lines.append("---")
            lines.append("")

        md_content = "\n".join(lines)
        file_path.write_text(md_content, encoding="utf-8")

        file_size = file_path.stat().st_size
        rel = _rel_path(file_path)

        return (
            f"Markdown document generated successfully.\n"
            f"Title: {title}\n"
            f"Sections: {len(sections)}\n"
            f"File: {rel}\n"
            f"Size: {file_size:,} bytes\n"
            f"Path: {str(file_path)}"
        )

    except Exception as exc:
        logger.error("generate_markdown failed: %s", exc)
        return f"Markdown generation failed: {exc}"


def make_generate_markdown_tool() -> StructuredTool:
    return StructuredTool.from_function(
        func=_generate_markdown,
        name="generate_markdown",
        description=(
            "Generate a Markdown (.md) document with a title, auto-generated table of contents, "
            "and structured sections. Each section has a 'heading' and 'content'. "
            "Returns the file path of the saved .md file."
        ),
        args_schema=GenerateMarkdownInput,
        return_direct=False,
    )
